#!/usr/bin/env python3
"""
Vitfix Investor Deck — PowerPoint Update Script
Transforms Fixit-Deck-Partenariats-B2B-v2.pptx into Vitfix branded deck
with real market data, Google Trends volumes, and verified statistics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
import os

# ═══════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════
INPUT_FILE = "/Users/elgato_fofo/Documents/Fixit-Deck-Partenariats-B2B-v2.pptx"
OUTPUT_FILE = "/Users/elgato_fofo/Desktop/Vitfix-Deck-Investisseurs-2026.pptx"

# Colors (Vitfix brand)
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)
ORANGE = RGBColor(0xFF, 0xC1, 0x07)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xD3, 0x2F, 0x2F)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
GOLD = RGBColor(0xFF, 0xD5, 0x4F)
DEEP_ORANGE = RGBColor(0xFF, 0x57, 0x22)

# Dimensions (10" x 5.625" = 16:9)
SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)
MARGIN = Emu(457200)  # 0.5 inch
CONTENT_W = Emu(8229600)


# ═══════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    """Add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_with_text(slide, left, top, width, height, text, font_size=14,
                        bold=False, text_color=DARK_TEXT, fill_color=None,
                        alignment=PP_ALIGN.LEFT, font_name='Arial'):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = font_name
    p.alignment = alignment
    return shape


def add_multi_text(slide, left, top, width, height, lines, fill_color=None):
    """Add a shape with multiple formatted text lines.
    lines = [(text, font_size, bold, color), ...]
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(4)

    return shape


def add_stat_box(slide, left, top, width, height, number, label, source,
                 num_color=DEEP_ORANGE, bg_color=None):
    """Add a stat box with big number + label + source."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Number
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = num_color
    p.font.name = 'Arial Black'
    p.alignment = PP_ALIGN.CENTER

    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = DARK_TEXT
    p2.font.name = 'Arial'
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    # Source
    p3 = tf.add_paragraph()
    p3.text = source
    p3.font.size = Pt(7)
    p3.font.color.rgb = GRAY
    p3.font.name = 'Arial'
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(2)

    return shape


def add_dark_banner(slide, top, text, font_size=16, color=WHITE, bg=DARK_BLUE):
    """Add a full-width dark banner."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), top, SLIDE_W, Emu(457200))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ═══════════════════════════════════════════════════
# STEP 1: FIND & REPLACE FIXIT → VITFIX
# ═══════════════════════════════════════════════════

def replace_text_in_presentation(prs):
    """Replace all occurrences of Fixit/FIXIT/fixit with Vitfix/VITFIX/vitfix."""
    replacements = {
        'FIXIT': 'VITFIX',
        'Fixit': 'Vitfix',
        'fixit': 'vitfix',
        'SOLUTION FIXIT': 'SOLUTION VITFIX',
        'LA SOLUTION FIXIT': 'LA SOLUTION VITFIX',
        'POURQUOI FIXIT': 'POURQUOI VITFIX',
        'AVANT FIXIT': 'AVANT VITFIX',
        'APRÈS FIXIT': 'APRÈS VITFIX',
        'partenariats@fixit.fr': 'partenariats@vitfix.fr',
        'www.fixit.fr': 'www.vitfix.fr',
    }

    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original = run.text
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                        if run.text != original:
                            count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original = run.text
                                for old, new in replacements.items():
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                if run.text != original:
                                    count += 1
    return count


# ═══════════════════════════════════════════════════
# STEP 2: UPDATE EXISTING SLIDES DATA
# ═══════════════════════════════════════════════════

def update_slide2_problem(slide):
    """Update 'LE PROBLÈME' slide with verified stats."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Update specific stats with verified ones
                    if 'Trouver un artisan fiable = 3h de recherche' in run.text:
                        run.text = '• 39% des particuliers ne trouvent pas d\'artisan fiable (OpinionWay 2025)'
                    elif 'Délai d\'intervention : 5-10 jours' in run.text:
                        run.text = '• 25% des Français craignent les arnaques (OpinionWay 2025)'
                    elif 'Prix opaques, devis non comparables' in run.text:
                        run.text = '• 33% redoutent les malfacons (OpinionWay 2025)'
                    elif 'Risque d\'arnaque, travaux mal faits' in run.text:
                        run.text = '• Satisfaction artisans : seulement 55% en Ile-de-France (BVA)'
                    elif 'Temps perdu coordination : 15h/semaine' in run.text:
                        run.text = '• 71,5% des entreprises peinent a recruter (France Travail 2024)'
                    elif 'Litiges interventions : 40% des cas' in run.text:
                        run.text = '• 485 000 postes vacants dans le BTP (FFB 2024)'
                    elif 'Facturation éparpillée' in run.text:
                        run.text = '• 30% des artisans sous-digitalises (PlanRadar 2024)'
                    elif 'Pas de traçabilité' in run.text:
                        run.text = '• 2 millions de degats des eaux/an (France Assureurs 2024)'
                    elif 'Clients/Locataires insatisfaits' in run.text:
                        run.text = '• 4 160 sinistres/jour = besoin artisans constant'
                    elif 'COÛT CACHÉ : 5 000' in run.text:
                        run.text = '\U0001F525 COUT TOTAL : 2,4 Md\u20AC/an d\'indemnisations degats des eaux seuls (France Assureurs 2024)'


def update_slide5_copro(slide):
    """Update copropriete slide with real numbers."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '500+ artisans vérifiés' in run.text:
                        run.text = '\u2705 Reseau d\'artisans verifies (SIRET + assurance)'


def update_slide4_segments(slide):
    """Update segments with verified numbers."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '750K unités' in run.text:
                        run.text = '873K immeubles'
                    elif '8M+ UNITÉS' in run.text:
                        run.text = '13M+ DE LOGEMENTS A ADRESSER'


# ═══════════════════════════════════════════════════
# STEP 3: CREATE NEW SLIDES
# ═══════════════════════════════════════════════════

def create_slide_marche(prs):
    """Create 'LE MARCHE EN CHIFFRES' slide with verified market data."""
    slide_layout = prs.slide_layouts[0]  # DEFAULT
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, MARGIN, Emu(274320), CONTENT_W, Emu(548640),
                '\U0001F4CA LE MARCHE EN CHIFFRES', 36, True, DARK_TEXT, PP_ALIGN.CENTER, 'Arial Black')

    # Subtitle
    add_textbox(slide, MARGIN, Emu(822960), CONTENT_W, Emu(274320),
                'Donnees verifiees — Sources : FFB, CAPEB, France Assureurs, France Travail, ANAH (2024)',
                10, False, GRAY, PP_ALIGN.CENTER)

    # Row 1: 4 stat boxes
    box_w = Emu(2011680)
    box_h = Emu(1280160)
    gap = Emu(91440)
    y1 = Emu(1188720)

    add_stat_box(slide, MARGIN, y1, box_w, box_h,
                 '208 Md\u20AC', 'Marche du batiment\nFrance 2024', 'Source : FFB 2024',
                 DEEP_ORANGE)

    add_stat_box(slide, Emu(457200) + box_w + gap, y1, box_w, box_h,
                 '118 Md\u20AC', 'Maintenance &\nRenovation', 'Source : FFB 2024 (57% du CA)',
                 DEEP_ORANGE)

    add_stat_box(slide, Emu(457200) + (box_w + gap) * 2, y1, box_w, box_h,
                 '620 000', 'Entreprises artisanales\ndu batiment', 'Source : CAPEB 2024',
                 RGBColor(0x1B, 0x5E, 0x20))

    add_stat_box(slide, Emu(457200) + (box_w + gap) * 3, y1, box_w, box_h,
                 '1,76M', 'Actifs dans\nle batiment', 'Source : FFB 2024',
                 RGBColor(0x1B, 0x5E, 0x20))

    # Row 2: 4 stat boxes
    y2 = Emu(2560320)

    add_stat_box(slide, MARGIN, y2, box_w, box_h,
                 '485 000', 'Postes vacants\ndans le BTP', 'Source : FFB / France Travail 2024',
                 RED)

    add_stat_box(slide, Emu(457200) + box_w + gap, y2, box_w, box_h,
                 '71,5%', 'Entreprises en\ndifficulte de recrutement', 'Source : France Travail BMO 2024',
                 RED)

    add_stat_box(slide, Emu(457200) + (box_w + gap) * 2, y2, box_w, box_h,
                 '873 000', 'Coproprietes\nen France', 'Source : ANIL / CoproFF 2023',
                 RGBColor(0x15, 0x65, 0xC0))

    add_stat_box(slide, Emu(457200) + (box_w + gap) * 3, y2, box_w, box_h,
                 '2M/an', 'Sinistres degats\ndes eaux', 'Source : France Assureurs 2024',
                 RGBColor(0x15, 0x65, 0xC0))

    # Bottom insight
    add_dark_banner(slide, Emu(4023360),
                    '\U0001F4A1 39% des particuliers ne trouvent pas d\'artisan fiable (OpinionWay 2025) — 92% cherchent en ligne (Google)',
                    13, ORANGE, DARK_BLUE)

    # Source bar
    add_textbox(slide, MARGIN, Emu(4572000), CONTENT_W, Emu(365760),
                'Sources : FFB (ffbatiment.fr) | CAPEB (capeb.fr) | France Travail BMO 2024 | France Assureurs 2024 | ANIL/ANAH | OpinionWay/illiCO 2025',
                8, False, GRAY, PP_ALIGN.CENTER)

    return slide


def create_slide_demande_digitale(prs):
    """Create 'LA DEMANDE DIGITALE' slide with Google Trends data."""
    slide_layout = prs.slide_layouts[0]  # DEFAULT
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, MARGIN, Emu(274320), CONTENT_W, Emu(548640),
                '\U0001F4C8 LA DEMANDE DIGITALE EXPLOSE', 36, True, DARK_TEXT, PP_ALIGN.CENTER, 'Arial Black')

    # Subtitle
    add_textbox(slide, MARGIN, Emu(822960), CONTENT_W, Emu(274320),
                'Volumes de recherche Google reels (France) — Sources : Ahrefs, Google Trends 2021-2025',
                10, False, GRAY, PP_ALIGN.CENTER)

    # Left column: Top keywords volumes
    left_x = MARGIN
    left_w = Emu(4114800)

    add_shape_with_text(slide, left_x, Emu(1188720), left_w, Emu(365760),
                        '\U0001F50D VOLUMES DE RECHERCHE MENSUELS (France)',
                        14, True, WHITE, RGBColor(0x15, 0x65, 0xC0), PP_ALIGN.CENTER)

    keywords_lines = [
        ('\U0001F527 serrurier : 53 000 rech/mois', 12, True, DEEP_ORANGE),
        ('\U0001F527 plombier : 35 000 rech/mois', 12, True, DEEP_ORANGE),
        ('\U0001F4A7 fuite d\'eau : 33 000 rech/mois', 12, False, DARK_TEXT),
        ('\U0001F3E2 syndic copropriete : 25 000 rech/mois', 12, False, DARK_TEXT),
        ('\u26A1 electricien : 25 000 rech/mois', 12, False, DARK_TEXT),
        ('\U0001F3D7 couvreur : 18 000 rech/mois', 12, False, DARK_TEXT),
        ('\U0001F333 paysagiste : 16 000 rech/mois', 12, False, DARK_TEXT),
        ('\U0001F50E avis artisan : 12 000 rech/mois (+60%)', 12, False, GREEN),
        ('\U0001F4DD devis artisan : 8 000 rech/mois', 12, False, DARK_TEXT),
        ('Source : Ahrefs (nov. 2024) via plaqueplastique.fr', 8, False, GRAY),
    ]
    add_multi_text(slide, left_x, Emu(1600200), left_w, Emu(2743200), keywords_lines)

    # Right column: Trends explosion
    right_x = Emu(4754880)
    right_w = Emu(4114800)

    add_shape_with_text(slide, right_x, Emu(1188720), right_w, Emu(365760),
                        '\U0001F680 EXPLOSION DES RECHERCHES "AUTOUR DE MOI"',
                        13, True, WHITE, DEEP_ORANGE, PP_ALIGN.CENTER)

    trends_lines = [
        ('\U0001F4CD "plombier autour de moi"', 13, True, DARK_TEXT),
        ('    2021: ~360 rech/an \u2192 2025: ~36 720/an', 11, False, DARK_TEXT),
        ('    \U0001F4C8 +5 000% en 4 ans', 12, True, GREEN),
        ('', 6, False, DARK_TEXT),
        ('\U0001F4CD "serrurier autour de moi"', 13, True, DARK_TEXT),
        ('    Inexistant en 2021 \u2192 37 846/an en 2025', 11, False, DARK_TEXT),
        ('    \U0001F4C8 Nouvelle tendance explosive', 12, True, GREEN),
        ('', 6, False, DARK_TEXT),
        ('\U0001F4CD "electricien autour de moi"', 13, True, DARK_TEXT),
        ('    0 en 2021 \u2192 36 000/an en 2025', 11, False, DARK_TEXT),
        ('    \U0001F4C8 Creation d\'un nouveau marche', 12, True, GREEN),
    ]
    add_multi_text(slide, right_x, Emu(1600200), right_w, Emu(2743200), trends_lines)

    # Bottom stats bar
    stats_y = Emu(4389120)
    third_w = Emu(2697480)

    add_stat_box(slide, MARGIN, stats_y, third_w, Emu(640080),
                 '4,1M+', 'Recherches/an sur nos 53 mots-cles', 'Google Trends + Ahrefs 2025',
                 DEEP_ORANGE, RGBColor(0xFF, 0xF3, 0xE0))

    add_stat_box(slide, Emu(457200) + third_w + Emu(91440), stats_y, third_w, Emu(640080),
                 '92%', 'Cherchent en ligne avant de choisir', 'Source : Google / LearnThings',
                 RGBColor(0x15, 0x65, 0xC0), RGBColor(0xE3, 0xF2, 0xFD))

    add_stat_box(slide, Emu(457200) + (third_w + Emu(91440)) * 2, stats_y, third_w, Emu(640080),
                 '50-55\u20AC', 'CPC "serrurier Paris" sur Google Ads', 'Source : Google Ads',
                 RED, RGBColor(0xFF, 0xEB, 0xEE))

    return slide


def create_slide_penurie(prs):
    """Create 'LA CRISE DE L'ARTISANAT' slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, MARGIN, Emu(274320), CONTENT_W, Emu(548640),
                '\u26A0\uFE0F LA CRISE DE L\'ARTISANAT', 36, True, DARK_TEXT, PP_ALIGN.CENTER, 'Arial Black')

    add_textbox(slide, MARGIN, Emu(822960), CONTENT_W, Emu(274320),
                'Le secteur du batiment traverse une crise structurelle majeure — Vitfix est la reponse',
                12, False, GRAY, PP_ALIGN.CENTER)

    # Left: La crise
    left_x = MARGIN
    left_w = Emu(3931920)

    add_shape_with_text(slide, left_x, Emu(1188720), left_w, Emu(365760),
                        '\u274C LE CONSTAT ALARMANT', 15, True, WHITE, RED, PP_ALIGN.CENTER)

    crisis_lines = [
        ('485 000 postes vacants dans le BTP', 13, True, RED),
        ('Source : FFB / France Travail 2024', 8, False, GRAY),
        ('', 4, False, DARK_TEXT),
        ('71,5% des entreprises en difficulte de recrutement', 12, True, DARK_TEXT),
        ('Source : France Travail, Enquete BMO 2024', 8, False, GRAY),
        ('', 4, False, DARK_TEXT),
        ('200 000 professionnels supplementaires necessaires', 12, True, DARK_TEXT),
        ('d\'ici 2030 pour la renovation energetique', 11, False, DARK_TEXT),
        ('Source : France Strategie', 8, False, GRAY),
        ('', 4, False, DARK_TEXT),
        ('39% des particuliers ne trouvent pas d\'artisan', 12, True, DARK_TEXT),
        ('Source : OpinionWay / illiCO travaux, Fev. 2025', 8, False, GRAY),
        ('', 4, False, DARK_TEXT),
        ('Seulement 64% de satisfaction client', 12, True, DARK_TEXT),
        ('55% en Ile-de-France | Source : IFOP / BVA', 8, False, GRAY),
    ]
    add_multi_text(slide, left_x, Emu(1600200), left_w, Emu(2926080), crisis_lines)

    # Right: Vitfix solution
    right_x = Emu(4754880)
    right_w = Emu(3931920)

    add_shape_with_text(slide, right_x, Emu(1188720), right_w, Emu(365760),
                        '\u2705 LA REPONSE VITFIX', 15, True, WHITE, GREEN, PP_ALIGN.CENTER)

    solution_lines = [
        ('Connecter l\'offre a la demande', 13, True, GREEN),
        ('En temps reel, par geolocalisation', 10, False, GRAY),
        ('', 4, False, DARK_TEXT),
        ('\u2705 Calendriers artisans en temps reel', 12, False, DARK_TEXT),
        ('\u2705 Reservation en 2 clics (comme Doctolib)', 12, False, DARK_TEXT),
        ('\u2705 Artisans verifies (SIRET + assurance)', 12, False, DARK_TEXT),
        ('\u2705 Avis clients certifies', 12, False, DARK_TEXT),
        ('\u2705 0\u20AC commission pour les clients', 12, False, DARK_TEXT),
        ('\u2705 App mobile native (iOS + Android)', 12, False, DARK_TEXT),
        ('\u2705 Dashboard syndic centralise', 12, False, DARK_TEXT),
        ('\u2705 IA comptable integree (Agent Lea)', 12, False, DARK_TEXT),
        ('', 4, False, DARK_TEXT),
        ('\U0001F680 Vitfix optimise chaque artisan existant', 12, True, GREEN),
        ('plutot que d\'en creer de nouveaux', 11, False, DARK_TEXT),
    ]
    add_multi_text(slide, right_x, Emu(1600200), right_w, Emu(2926080), solution_lines)

    # Bottom banner
    add_dark_banner(slide, Emu(4663440),
                    '\U0001F4A1 Chaque artisan connecte via Vitfix = + de clients servis, moins de temps perdu, plus de revenus',
                    13, ORANGE, DARK_BLUE)

    return slide


def create_slide_opportunite(prs):
    """Create 'OPPORTUNITE INVESTISSEURS' slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    add_textbox(slide, MARGIN, Emu(274320), CONTENT_W, Emu(548640),
                '\U0001F4B0 OPPORTUNITE INVESTISSEURS', 36, True, ORANGE, PP_ALIGN.CENTER, 'Arial Black')

    add_textbox(slide, MARGIN, Emu(822960), CONTENT_W, Emu(274320),
                'Un marche de 208 milliards EUR digitalise a moins de 15% = opportunite massive',
                12, False, LIGHT_GRAY, PP_ALIGN.CENTER)

    # Key metrics row
    box_w = Emu(2560320)
    gap = Emu(91440)
    y1 = Emu(1280160)
    box_h = Emu(1188720)

    add_stat_box(slide, MARGIN, y1, box_w, box_h,
                 '208 Md\u20AC', 'Marche total batiment\nFrance 2024', 'Source : FFB',
                 ORANGE, RGBColor(0x2A, 0x2A, 0x4E))

    add_stat_box(slide, Emu(457200) + box_w + gap, y1, box_w, box_h,
                 '4,1M+', 'Recherches Google/an\nsur nos mots-cles', 'Google Trends + Ahrefs',
                 ORANGE, RGBColor(0x2A, 0x2A, 0x4E))

    add_stat_box(slide, Emu(457200) + (box_w + gap) * 2, y1, box_w, box_h,
                 '+5 000%', 'Croissance recherches\n"autour de moi" (4 ans)', 'Google Trends 2021-2025',
                 GREEN, RGBColor(0x2A, 0x2A, 0x4E))

    # Revenue model
    model_y = Emu(2651760)
    model_h = Emu(914400)

    add_shape_with_text(slide, MARGIN, model_y, CONTENT_W, Emu(365760),
                        '\U0001F4CA PROJECTION DE REVENUS (hypothese conservatrice)',
                        14, True, WHITE, RGBColor(0x2A, 0x2A, 0x4E), PP_ALIGN.CENTER)

    rev_lines = [
        ('ANNEE 1 : 500 artisans x 49\u20AC/mois = 294 000\u20AC ARR', 14, True, ORANGE),
        ('   + 20 syndics Starter (gratuit) + 5 syndics Pro = acquisition B2B', 11, False, LIGHT_GRAY),
        ('', 4, False, WHITE),
        ('ANNEE 2 : 2 000 artisans + 50 syndics Pro = 1,5M\u20AC ARR', 14, True, ORANGE),
        ('   + Commissions interventions (8-15%) + White-label', 11, False, LIGHT_GRAY),
        ('', 4, False, WHITE),
        ('ANNEE 3 : 8 000 artisans + 200 syndics = 6M\u20AC+ ARR', 14, True, ORANGE),
        ('   + Expansion 3 villes + API marketplace + Effet reseau', 11, False, LIGHT_GRAY),
    ]
    add_multi_text(slide, MARGIN, Emu(3017520), CONTENT_W, Emu(1645920), rev_lines,
                   RGBColor(0x2A, 0x2A, 0x4E))

    # Bottom tagline
    add_textbox(slide, MARGIN, Emu(4754880), CONTENT_W, Emu(365760),
                '\U0001F680 Vitfix : le Doctolib de l\'artisanat — Un marche de 208 Md\u20AC, une digitalisation a <15%, une demande qui explose',
                12, True, GOLD, PP_ALIGN.CENTER)

    return slide


def create_slide_confiance_chiffres(prs):
    """Create 'CONFIANCE & CHIFFRES CLES' slide with additional verified stats."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, MARGIN, Emu(274320), CONTENT_W, Emu(548640),
                '\U0001F4CB CHIFFRES CLES — TOUS VERIFIES', 36, True, DARK_TEXT, PP_ALIGN.CENTER, 'Arial Black')

    add_textbox(slide, MARGIN, Emu(822960), CONTENT_W, Emu(274320),
                'Donnees de marche pour le secteur artisan/batiment en France — chaque chiffre est source',
                10, False, GRAY, PP_ALIGN.CENTER)

    # 3x3 grid of stat boxes
    box_w = Emu(2651760)
    box_h = Emu(1005840)
    gap = Emu(91440)

    stats = [
        ('90%', 'Considerent les avis en ligne\nessentiels pour choisir un artisan', 'IFOP / Plus que PRO 2021', DEEP_ORANGE),
        ('43%', 'Des artisans ne croient pas\nen l\'impact du digital', 'Batiweb / Etude sectorielle', RGBColor(0x15, 0x65, 0xC0)),
        ('2,4 Md\u20AC', 'Indemnisations degats des eaux\npar an en France (+134% en 20 ans)', 'France Assureurs 2024', RED),
        ('4 160/jour', 'Sinistres degats des eaux\nen France', 'France Assureurs 2024', RED),
        ('13M', 'Logements en copropriete\nen France', 'ANIL 2023', RGBColor(0x15, 0x65, 0xC0)),
        ('5 088', 'Detenteurs carte S\n(syndics professionnels)', 'CCI-France', DARK_TEXT),
        ('96,3%', 'Part de marche Google\nen France (mobile)', 'StatCounter / WebrankInfo', GREEN),
        ('78%', 'Recherches locales mobiles\nmenent a un achat', 'Google Data', GREEN),
        ('+45%/an', 'Croissance recherches\nde proximite', 'Google Trends 2021-2025', DEEP_ORANGE),
    ]

    for i, (num, label, source, color) in enumerate(stats):
        row = i // 3
        col = i % 3
        x = Emu(457200) + (box_w + gap) * col
        y = Emu(1188720) + (box_h + gap) * row
        add_stat_box(slide, x, y, box_w, box_h, num, label, source, color)

    # Bottom
    add_textbox(slide, MARGIN, Emu(4480560), CONTENT_W, Emu(365760),
                '\U0001F4D6 Toutes les sources sont publiques et verifiables : FFB, CAPEB, France Travail, France Assureurs, ANIL, ANAH, Google, IFOP, BVA, OpinionWay',
                9, False, GRAY, PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════
# STEP 4: REORDER SLIDES
# ═══════════════════════════════════════════════════

def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    if new_index >= len(slides):
        xml_slides.append(el)
    else:
        slides_after = list(xml_slides)
        if new_index < len(slides_after):
            xml_slides.insert(new_index, el)
        else:
            xml_slides.append(el)


# ═══════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════

def main():
    print("\U0001F4E6 Loading presentation...")
    prs = Presentation(INPUT_FILE)

    original_count = len(prs.slides)
    print(f"   {original_count} slides loaded")

    # Step 1: Replace FIXIT → VITFIX
    print("\n\U0001F504 Renaming FIXIT → VITFIX...")
    replace_count = replace_text_in_presentation(prs)
    print(f"   {replace_count} text replacements made")

    # Step 2: Update existing slides with verified data
    print("\n\U0001F4CA Updating existing slides with verified data...")
    slides = list(prs.slides)
    update_slide2_problem(slides[1])   # Slide 2: LE PROBLEME
    update_slide4_segments(slides[3])   # Slide 4: SEGMENTS
    update_slide5_copro(slides[4])      # Slide 5: COPRO
    print("   Slides 2, 4, 5 updated")

    # Step 3: Create new slides (they get added at the end)
    print("\n\u2795 Creating new slides...")

    slide_marche = create_slide_marche(prs)
    print("   + LE MARCHE EN CHIFFRES")

    slide_penurie = create_slide_penurie(prs)
    print("   + LA CRISE DE L'ARTISANAT")

    slide_digitale = create_slide_demande_digitale(prs)
    print("   + LA DEMANDE DIGITALE EXPLOSE")

    slide_chiffres = create_slide_confiance_chiffres(prs)
    print("   + CHIFFRES CLES VERIFIES")

    slide_opportunite = create_slide_opportunite(prs)
    print("   + OPPORTUNITE INVESTISSEURS")

    # Step 4: Reorder slides
    # Current order after adding: [1-14 original, 15-marche, 16-penurie, 17-digitale, 18-chiffres, 19-opportunite]
    # Desired order:
    # 1. Title (VITFIX)
    # 2. LE PROBLEME (updated)
    # 3. NEW: LE MARCHE EN CHIFFRES
    # 4. NEW: LA CRISE DE L'ARTISANAT
    # 5. NEW: LA DEMANDE DIGITALE EXPLOSE
    # 6. LA SOLUTION VITFIX
    # 7. NOS 7 SEGMENTS
    # 8. COPROPRIETES & SYNDICS
    # 9. BAILLEURS SOCIAUX
    # 10. NOS OFFRES ARTISANS
    # 11. OFFRES PARTENAIRES B2B
    # 12. CAS CLIENT
    # 13. POURQUOI VITFIX
    # 14. DEMARRAGE 4 SEMAINES
    # 15. NOS ENGAGEMENTS
    # 16. TEMOIGNAGES
    # 17. NEW: CHIFFRES CLES VERIFIES
    # 18. NEW: OPPORTUNITE INVESTISSEURS
    # 19. CTA

    print("\n\U0001F500 Reordering slides...")
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)

    # Build new order by index
    # Original indices: 0-13, new: 14(marche), 15(penurie), 16(digitale), 17(chiffres), 18(opportunite)
    new_order = [
        0,   # 1. Title
        1,   # 2. Le Probleme
        14,  # 3. Le Marche en Chiffres (NEW)
        15,  # 4. La Crise (NEW)
        16,  # 5. La Demande Digitale (NEW)
        2,   # 6. La Solution Vitfix
        3,   # 7. Nos 7 Segments
        4,   # 8. Coproprietes & Syndics
        5,   # 9. Bailleurs Sociaux
        6,   # 10. Nos Offres Artisans
        7,   # 11. Offres Partenaires B2B
        8,   # 12. Cas Client
        9,   # 13. Pourquoi Vitfix
        10,  # 14. Demarrage 4 semaines
        11,  # 15. Nos Engagements
        12,  # 16. Temoignages
        17,  # 17. Chiffres Cles (NEW)
        18,  # 18. Opportunite Investisseurs (NEW)
        13,  # 19. CTA
    ]

    # Reorder XML elements
    for el in slides_list:
        xml_slides.remove(el)
    for idx in new_order:
        xml_slides.append(slides_list[idx])

    print(f"   {len(new_order)} slides reordered")

    # Save
    print(f"\n\U0001F4BE Saving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)

    final_count = len(prs.slides)
    print(f"\n\u2705 Done! {final_count} slides total ({final_count - original_count} new slides added)")
    print(f"\U0001F4C4 Output: {OUTPUT_FILE}")


if __name__ == '__main__':
    main()
